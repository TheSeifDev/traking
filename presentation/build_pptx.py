# -*- coding: utf-8 -*-
"""Builds the Phantoms kickoff meeting deck (.pptx) in the team brand:
black + crimson-red + white, right-to-left Arabic."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
A = os.path.join(HERE, "assets")

# ---------- Brand palette ----------
BLACK   = RGBColor(0x08,0x08,0x0A)
PANEL   = RGBColor(0x15,0x15,0x19)
PANEL2  = RGBColor(0x1C,0x10,0x12)   # faint dark-red panel
RED     = RGBColor(0xE1,0x0B,0x06)
RED_D   = RGBColor(0x8E,0x07,0x04)
RED_GHOST=RGBColor(0xB815,0x10,0x10) if False else RGBColor(0xB8,0x15,0x10)
WHITE   = RGBColor(0xF4,0xF4,0xF7)
GREY    = RGBColor(0xC6,0xC6,0xCF)
DIM     = RGBColor(0x8C,0x8C,0x99)
LINE    = RGBColor(0x2E,0x2E,0x36)
LINE_R  = RGBColor(0x6E,0x14,0x10)
GREEN   = RGBColor(0x3D,0xDC,0x97)
AMBER   = RGBColor(0xFF,0xC2,0x4B)

AR_FONT = "Cairo"
LAT_FONT= "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

# ---------------- helpers ----------------
def slide():
    return prs.slides.add_slide(BLANK)

def _set_font(run, size, color, bold, font=AR_FONT, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin","a:cs","a:ea"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", font)
    # ensure complex-script (Arabic) rendering + bidi
    rPr.set("rtl","1")

def bg(s, glow=True):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BLACK; r.line.fill.background()
    r.shadow.inherit = False
    if glow:
        g = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.6), Inches(-2.6), Inches(6.2), Inches(6.2))
        g.fill.solid(); g.fill.fore_color.rgb = RGBColor(0x2A,0x06,0x05)
        g.line.fill.background(); g.shadow.inherit=False
        # soft transparent
        sp = g.fill._xPr.find(qn('a:solidFill'))
        clr = sp.find(qn('a:srgbClr')); a = clr.makeelement(qn('a:alpha'),{'val':'38000'}); clr.append(a)
    return r

def _para(tf, runs, align, space_after, space_before, line, first, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_after is not None: p.space_after = Pt(space_after)
    if space_before is not None: p.space_before = Pt(space_before)
    if line is not None:
        p.line_spacing = line
    # RTL
    pPr = p._p.get_or_add_pPr()
    pPr.set("rtl","1")
    for (txt, kw) in runs:
        run = p.add_run(); run.text = txt
        _set_font(run, kw.get("size",18), kw.get("color",GREY), kw.get("bold",False),
                  kw.get("font",AR_FONT), kw.get("italic",False))
    return p

def text(s, x, y, w, h, paras, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    first=True
    for pd in paras:
        if isinstance(pd, tuple): pd=[pd]
        _para(tf, pd, pd[0][1].get("align",align) if pd else align,
              None,None,None,first)
        # apply per-paragraph spacing/line carried in first run kwargs
        kw0 = pd[0][1]
        p = tf.paragraphs[-1]
        if kw0.get("after") is not None: p.space_after=Pt(kw0["after"])
        if kw0.get("before") is not None: p.space_before=Pt(kw0["before"])
        if kw0.get("line") is not None: p.line_spacing=kw0["line"]
        first=False
    return tb

def R(t,**kw):  # run shorthand
    return (t, kw)

def shape(s, kind, x,y,w,h, fill=None, line=None, line_w=1.0, shadow=False, radius=None):
    sp = s.shapes.add_shape(kind, Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else:
        sp.line.color.rgb=line; sp.line.width=Pt(line_w)
    sp.shadow.inherit = shadow
    if radius is not None and kind==MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0]=radius
        except Exception: pass
    return sp

def card(s, x,y,w,h, accent=False, fill=PANEL, line=LINE):
    if accent: fill=PANEL2; line=LINE_R
    sp = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x,y,w,h, fill=fill, line=line, line_w=1.25, radius=0.06)
    return sp

def kicker(s, x, y, txt, color=RED, w=None):
    w = w or (0.28 + 0.115*len(txt))
    sp = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.46, fill=None,
               line=color, line_w=1.25, radius=0.5)
    tf=sp.text_frame; tf.word_wrap=False; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_top=0;tf.margin_bottom=0;tf.margin_left=Inches(0.12);tf.margin_right=Inches(0.12)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; p._p.get_or_add_pPr().set("rtl","1")
    r=p.add_run(); r.text=txt; _set_font(r,13.5,color,True)
    return sp

def heading(s, x, y, w, runs, size=33, after=0):
    text(s, x, y, w, 1.1, [ [R(t, size=size, color=c, bold=True, after=after, line=1.15)] for (t,c) in runs ])

def logo(s, x, y, h):
    return s.shapes.add_picture(os.path.join(A,"logo.png"), Inches(x), Inches(y), height=Inches(h))

def footer(s, n):
    shape(s, MSO_SHAPE.RECTANGLE, 0.6, 7.06, 12.13, 0.012, fill=LINE)
    text(s, 0.6, 7.12, 5.0, 0.3, [[R("PHANTOMS", size=10, color=RED, bold=True, font=LAT_FONT),
                                   R("  —  MEETING 01 · 2026", size=10, color=DIM, font=LAT_FONT)]],
         align=PP_ALIGN.LEFT)
    text(s, 11.2, 7.12, 1.5, 0.3, [[R(f"{n:02d}", size=11, color=WHITE, bold=True, font=LAT_FONT),
                                    R(" / 18", size=10, color=DIM, font=LAT_FONT)]], align=PP_ALIGN.RIGHT)

def bullets(s, x,y,w,h, items, size=15.5, gap=8, color=GREY, marker_color=RED, line=1.3):
    paras=[]
    for it in items:
        if isinstance(it, str): it=[R(it,size=size,color=color)]
        else: it=[R("◈  ", size=size, color=marker_color, bold=True)]+it
        paras.append(it)
    text(s,x,y,w,h,paras, align=PP_ALIGN.RIGHT)

def pic_cover(path, x,y,w,h):
    """cover-crop an image into a w x h box."""
    iw,ih = Image.open(path).size
    tr = w/h; ir = iw/ih
    pic = s_pic = None
    if ir > tr:
        pic_h = h; pic_w = h*ir
    else:
        pic_w = w; pic_h = w/ir
    px = x - (pic_w-w)/2; py = y - (pic_h-h)/2
    return None  # placeholder, replaced below

def add_pic_crop(s, path, x,y,w,h, line=None):
    iw,ih = Image.open(path).size
    box_tr = w/h; img_ar = iw/ih
    if img_ar > box_tr:
        new_h=h; new_w=h*img_ar
    else:
        new_w=w; new_h=w/img_ar
    px = x - (new_w-w)/2; py = y - (new_h-h)/2
    pic = s.shapes.add_picture(path, Inches(px), Inches(py), Inches(new_w), Inches(new_h))
    # crop via crop factors
    if img_ar > box_tr:
        crop = (new_w-w)/2/new_w
        pic.crop_left=crop; pic.crop_right=crop
    else:
        crop = (new_h-h)/2/new_h
        pic.crop_top=crop; pic.crop_bottom=crop
    if line:
        pic.line.color.rgb=line; pic.line.width=Pt(1.5)
    return pic

# ============================================================ 1 COVER
s = slide(); bg(s)
logo(s, 0.9, 1.0, 4.6)
shape(s, MSO_SHAPE.RECTANGLE, 0.0, 5.02, SW, 0.03, fill=RED)
kicker(s, 7.0, 1.15, "أول ميتنج رسمي في السنة", w=3.2)
text(s, 6.2, 1.75, 6.4, 2.2, [
    [R("خطة ميتنج ", size=54, color=WHITE, bold=True, line=1.05),
     R("Phantoms", size=54, color=RED, bold=True, font=LAT_FONT)],
])
text(s, 6.2, 3.35, 6.5, 1.5, [
    [R("مش تيم بيحاول يبان كويس…", size=20, color=GREY, bold=True, after=4, line=1.4)],
    [R("تيم بيتبني منتجات حقيقية، وناس حقيقية.", size=20, color=WHITE, bold=True, line=1.4)],
])
for i,(t1,t2) in enumerate([("⏱  ساعتان كاملتان","10:00 — 12:00"),]):
    pass
chips = [("⏱  ساعتان كاملتان",),("🕙  10:00 — 12:00",),("🎯  استحواذ · تأسيس · تحوّل",)]
cx=6.2
for (t,) in chips:
    wch = 0.55 + 0.128*len(t)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, 5.35, wch, 0.6, fill=PANEL, line=LINE, line_w=1, radius=0.5)
    text(s, cx, 5.46, wch, 0.4, [[R(t, size=14.5, color=WHITE, bold=True)]], align=PP_ALIGN.CENTER)
    cx += wch + 0.25

# ============================================================ 2 TIMELINE
s=slide(); bg(s); footer(s,2)
kicker(s, 10.0, 0.55, "الجدول الزمني الواقعي", w=3.0)
heading(s, 0.6, 1.05, 12.1, [("الساعتان، ",WHITE),("دقيقة بدقيقة",RED)], size=34)
rows = [
 ("10:00 – 10:10","مرحلة الانتظار — 10 دقائق","حضور وتسجيل وتهيئة، من غير بداية رسمية قبل الميعاد.", AMBER),
 ("10:10 – 10:25","الافتتاحية: الهوية + ليه بدأنا بدري — 15 دقيقة","ترحيب، كسر فخ الاعتذار، تعريف مختصر، الجملة التأسيسية.", RED),
 ("10:25 – 10:35","ليه من غير فكرة محددة + البرمشن — 10 دقائق","قانون الأساس قبل الفكرة، ونقطة التفويض الرسمي — جملة واحدة واضحة.", RED),
 ("10:35 – 11:15","عروض اللجان (Heads + Monitors) — 40 دقيقة","3 لجان × 12–13 دقيقة: رؤية، تقييم، Roadmap، أدوار، معايير استمرار.", RED),
 ("11:15 – 11:25","الصندوق الأسود — أسئلة الفورم — 10 دقائق","رد سريع ومباشر على أكتر الأسئلة تكرارًا، من غير ذكر أسماء.", RED),
 ("11:25 – 11:55","Mindset + مهارتين بعمق — 30 دقيقة","3 رسائل جوهرية بس، كل واحدة بعمق: عقلية، 80/20، تواصل بحضور.", RGBColor(0xB8,0x15,0x10)),
 ("11:55 – 12:00","الختام والتحوّل — 5 دقائق","خلاصة واحدة، تقدير للحضور… وبعدها اللعب يبدأ فورًا بلا فاصل.", AMBER),
]
y=1.95
for (tm,ti,de,col) in rows:
    shape(s, MSO_SHAPE.OVAL, 2.55, y+0.12, 0.16,0.16, fill=BLACK, line=col, line_w=2.2)
    text(s, 0.6, y, 1.85, 0.5, [[R(tm, size=13.5, color=col, bold=True, font=LAT_FONT)]], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    card(s, 2.95, y-0.06, 9.75, 0.62, fill=PANEL, line=LINE)
    text(s, 3.15, y, 9.4, 0.6, [
        [R(ti+"   ", size=14.5, color=WHITE, bold=True), R(de, size=12.5, color=GREY)]
    ], anchor=MSO_ANCHOR.MIDDLE)
    y+=0.71

# ============================================================ 3 OPENING
s=slide(); bg(s); footer(s,3)
kicker(s, 10.7, 0.55, "الافتتاحية · 15 دقيقة", w=2.4)
heading(s, 0.6, 1.05, 12.1, [("أول 15 دقيقة هدفها ",WHITE),("استحواذ الانتباه",RED),("، مش سرد إنجازات",WHITE)], size=29)
cards3 = [
 (0.6,1.95,"🗣  جملة الدخول — وكسر فخ الاعتذار",
   [[R("«صباح الخير جميعًا… شكرًا لالتزامكم وحضوركم في الميعاد، وشكرًا لانتظاركم.»",size=14,color=WHITE,bold=True,line=1.4,after=6)],
    [R("شكرًا على الالتزام — ",size=13,color=GREY),R("مش",size=13,color=RED,bold=True),R(" «آسف على التأخير». أول جملة بتحدد نبرة القائد.",size=13,color=GREY,line=1.4)]], True),
 (6.95,1.95,"🪪  مين أنا — في جملتين بس",
   [[R("خبرتك التقنية + إنجاز مختصر بأرقام يخلق مصداقية بسرعة، وبعدها فورًا:",size=13,color=GREY,line=1.4,after=6)],
    [R("«أنا مش هنا أقولكم عملت إيه… أنا هنا أوفّرلكم البيئة اللي تخليكم تعملوا أحسن من كده.»",size=13.5,color=WHITE,bold=True,line=1.4,after=6)],
    [R("⚠ مفيش قصة حياة. دقيقتين بالظبط.",size=12.5,color=AMBER,bold=True)]], False),
 (0.6,4.35,"👻  هوية Phantoms",
   [[R("«إحنا مش تيم بيحاول يبان كويس… إحنا تيم بيتبني ",size=14,color=GREY,line=1.45),R("منتجات حقيقية وناس حقيقية",size=14,color=WHITE,bold=True),R(".»",size=14,color=GREY)]], False),
 (6.95,4.35,"⚓  الجملة التأسيسية — شعار السنة",
   [[R("جملة واحدة قصيرة تتكرر طول السنة وتبقى علامة التيم. بتتعرض في البداية، وتتكرر في الختام، وتبقى على كل حاجة بعد كده.",size=14,color=GREY,line=1.45)]], True),
]
for (x,y,t,paras,acc) in cards3:
    card(s,x,y,5.78,2.18,accent=acc)
    text(s,x+0.3,y+0.22,5.25,0.5,[[R(t,size=16.5,color=RED if acc else WHITE,bold=True)]])
    text(s,x+0.3,y+0.78,5.25,1.3,paras)

# ============================================================ 4 WHY EARLY
s=slide(); bg(s); footer(s,4)
kicker(s, 10.9, 0.55, "ليه بدأنا بدري · 5 دقائق", w=2.2)
heading(s, 0.6, 1.05, 12.1, [("بنشتري وقت… ",WHITE),("مش بنسرقه من حد",RED)], size=34)
card(s,0.6,2.0,5.9,2.7, fill=RGBColor(0x24,0x10,0x10), line=RGBColor(0x6E,0x14,0x10))
text(s,0.9,2.25,5.3,0.5,[[R("😵  اللي يبدأ مع أول محاضرة",size=17,color=RED,bold=True)]])
bullets(s,0.9,2.85,5.35,1.7,[
  [R("الضغط الأكاديمي بيسرق فرصة التأسيس الحقيقي.",size=14,color=GREY)],
  [R("طول الترم في وضع ",size=14,color=GREY),R("رد الفعل — Reactive",size=14,color=WHITE,bold=True),R(": بيلحق بس، بيتعلم في نص طاقة.",size=14,color=GREY)],
], size=14, gap=8)
card(s,6.85,2.0,5.9,2.7, fill=RGBColor(0x0E,0x1C,0x16), line=RGBColor(0x14,0x4A,0x36))
text(s,7.15,2.25,5.3,0.5,[[R("🎯  إحنا النهاردة",size=17,color=GREEN,bold=True)]])
bullets(s,7.15,2.85,5.35,1.7,[
  [R("بنبدأ بدري عشان نكون في وضع ",size=14,color=GREY),R("المبادرة — Proactive",size=14,color=WHITE,bold=True),R(".",size=14,color=GREY)],
  [R("إحنا اللي بنتحكم في الوتيرة، وبنبني الأساس بهدوء وعمق قبل ما الزحمة تبدأ.",size=14,color=GREY)],
], size=14, marker_color=GREEN, gap=8)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 5.15, 12.15, 1.25, fill=PANEL2, line=LINE_R, line_w=1.25, radius=0.08)
text(s,0.95,5.4,11.5,0.9,[[
  R("💡 الأساسيات اللي بنشتغل عليها دلوقتي ",size=15,color=WHITE,bold=True),
  R("مش حمل زيادة",size=15,color=RED,bold=True),
  R(" — هي أصلًا جزء من منهج السنة دي، بس بندرسها بعمق وتطبيق عملي بدل الحفظ السطحي وقت الامتحان.",size=15,color=GREY,line=1.5)]])

# ============================================================ 5 NO FIXED IDEA
s=slide(); bg(s); footer(s,5)
kicker(s, 9.35, 0.55, "ليه من غير فكرة محددة · 7 دقائق — Product Thinking", w=3.85, color=AMBER)
heading(s, 0.6, 1.05, 12.1, [("قانون ",WHITE),("الأساس قبل الفكرة",RED)], size=34)
card(s,0.6,1.95,12.13,1.55,accent=True)
text(s,0.95,2.2,11.5,1.2,[[R("«سواء اشتريت عربية مانيوال أو أوتوماتيك… لازم تكون متعلّم أساسيات السواقة الأول. لو حددنا فكرة دلوقتي، هنحبس نفسنا في سقف مهاراتنا الحالية. لما كل مسار — سوفت، هارد، AI — يمتلك أدواته بالكامل، هنقدر ننفّذ أي فكرة معقدة تكتسح أي منافسة.»",size=15.5,color=WHITE,bold=True,line=1.5)]])
three=[
 ("🧩  الفكرة مش هي القيمة",[[R("القيمة الحقيقية هي ",size=13.5,color=GREY),R("المهارة",size=13.5,color=WHITE,bold=True),R(" اللي هتنفّذ بيها. الأفكار رخيصة، التنفيذ هو اللي نادر.",size=13.5,color=GREY,line=1.4)]]),
 ("⚖️  القاعدة الذهبية",[[R("فكرة كويسة + فريق ضعيف = ",size=13.5,color=GREY),R("فشل",size=13.5,color=RED,bold=True),R(".",size=13.5,color=GREY)],[R("فريق قوي + أي فكرة = ",size=13.5,color=GREY),R("منتج ينافس",size=13.5,color=GREEN,bold=True),R(".",size=13.5,color=GREY,line=1.4)]]),
 ("🧠  عقلية Product",[[R("مش بتوع «مشروع تخرج يتسلّم»… بتوع منتج حقيقي بيتحلّ بيه مشكلة ناس حقيقية.",size=13.5,color=GREY,line=1.4)]]),
]
x=0.6
for (t,paras) in three:
    card(s,x,3.75,3.94,1.85,accent=True)
    text(s,x+0.28,3.98,3.4,0.5,[[R(t,size=15,color=WHITE,bold=True)]])
    text(s,x+0.28,4.5,3.45,1.0,paras)
    x+=4.09
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,0.6,5.85,12.13,1.05,fill=PANEL,line=LINE,radius=0.08)
text(s,0.95,6.05,11.5,0.8,[[R("🔎 بعدين لما نيجي نختار: ",size=13.5,color=AMBER,bold=True),
  R("(1) إيه المشكلة الحقيقية ولمين بالظبط؟  ·  (2) هل حد جاهز يدفع/يستخدم الحل فعلًا؟  ·  (3) إيه أصغر نسخة تجربة (MVP) قبل ما نبني كل حاجة؟ — التفاصيل في المستند المكتوب.",size=13.5,color=GREY,line=1.4)]])

# ============================================================ 6 PERMISSION
s=slide(); bg(s); footer(s,6)
kicker(s, 11.35, 0.55, "البرمشن · 8 دقائق", w=1.75, color=AMBER)
heading(s, 0.6, 1.05, 12.1, [("القوة في ",WHITE),("الوضوح",RED),("… مش في الغموض",WHITE)], size=32)
card(s,0.6,2.0,5.9,3.1, fill=RGBColor(0x24,0x10,0x10), line=RGBColor(0x6E,0x14,0x10))
text(s,0.9,2.25,5.3,0.5,[[R("🚫  اللي منعملوش",size=17,color=RED,bold=True)]])
bullets(s,0.9,2.9,5.35,2.1,[
 [R("منحوّلهاش «سر مرعب غامض» — الغموض الزيادة بيولّد ",size=14,color=GREY),R("قلق مش هيبة",size=14,color=WHITE,bold=True),R(".",size=14,color=GREY)],
 [R("منقولش «ممنوع تتكلم عن البرمشن دي» بنبرة تهديد.",size=14,color=GREY)],
 [R("منكررهاش أكتر من مرة — التكرار بيثير الفضول أكتر ما بيرسّخ الهيبة.",size=14,color=GREY)],
],gap=10)
card(s,6.85,2.0,5.9,3.1, fill=RGBColor(0x0E,0x1C,0x16), line=RGBColor(0x14,0x4A,0x36))
text(s,7.15,2.25,5.3,0.5,[[R("✅  الجملة — بحزم وهدوء، مرة واحدة",size=16,color=GREEN,bold=True)]])
text(s,7.15,2.9,5.35,2.1,[[R("«إحنا معانا تفويض رسمي إننا نبدأ ونتحرك من بدري. ده امتياز — والتفاصيل دي داخلية وخاصة بالتيم، وممنوع تخرج برا القاعة دي لأسباب تنظيمية، مش لأنه سر مخيف.»",size=15,color=WHITE,bold=True,line=1.55)]])
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,0.6,5.4,12.13,1.0,fill=PANEL2,line=LINE_R,line_w=1.25,radius=0.1)
text(s,0.95,5.62,11.5,0.7,[[R("🎯 نفس القوة، من غير ما الناس تحس إنك بتخبي حاجة وحشة. ",size=15,color=WHITE,bold=True),R("جملة واحدة واضحة… وكفاية.",size=15,color=GREY)]])

# ============================================================ 7 COMMITTEES
s=slide(); bg(s); footer(s,7)
kicker(s, 8.7, 0.55, "عروض اللجان · 40 دقيقة — Heads + Monitors", w=4.5)
heading(s, 0.6, 1.05, 12.1, [("3 لجان × ",WHITE),("12–13 دقيقة",RED),("… وكل لجنة تغطي 5 حاجات",WHITE)], size=27)
items=[
 ("1 · الرؤية الفنية","إيه مسار اللجنة (سوفت / هارد / AI)، ومين المفروض نوصل فين بنهاية الترم."),
 ("2 · تقييم الفترة السابقة","تقييم صريح لجاهزية الأعضاء — فين القوة وفين الفجوة."),
 ("3 · الـ Roadmap","خطة الترم الأول بتواريخ واضحة، مش عناوين."),
 ("4 · الأدوار","مين بيعمل إيه بالظبط — مفيش «هنعمل» من غير صاحب."),
 ("5 · معايير الاستمرار","نظام الـ Strikes والالتزام — واضح من الأول، من غير مفاجآت."),
 ("⏱  انضباط التوقيت","يُبلَّغوا قبلها بيوم إنك هتوقفهم لو زادوا. القائد واقف جنبهم — مش قاعد."),
]
x0,y0=0.6,1.95; cw,ch=3.94,1.55; gx,gy=0.19,0.2
for i,(t,d) in enumerate(items):
    r,c=divmod(i,3)
    x=x0+c*(cw+gx); y=y0+r*(ch+gy)
    acc = i<5
    card(s,x,y,cw,ch,accent=acc)
    text(s,x+0.26,y+0.18,cw-0.5,0.5,[[R(t,size=14.5,color=RED if acc else AMBER,bold=True)]])
    text(s,x+0.26,y+0.66,cw-0.5,0.85,[[R(d,size=12.5,color=GREY,line=1.3)]])
text(s,0.6,5.55,12.1,0.4,[[R("تسلسل التواصل الرسمي",size=15,color=WHITE,bold=True)]],align=PP_ALIGN.CENTER)
chain=[("Monitor","متابعة يومية + أعذار"),("Head","مشاكل فنية"),("سهيلة",""),("سيف","")]
bw=2.55; gap=0.35; total=len(chain)*bw+(len(chain)-1)*gap
cx=(13.333-total)/2
for i,(a,b) in enumerate(chain):
    x=cx+i*(bw+gap)
    head = (a=="سيف")
    card(s,x,6.0,bw,0.85, fill=PANEL2 if head else PANEL, line=LINE_R if head else LINE)
    text(s,x,6.12,bw,0.6,[[R(a,size=15,color=RED if head else WHITE,bold=True,font=(LAT_FONT if a in("Monitor","Head") else AR_FONT))],
                          [R(b,size=10.5,color=DIM,line=1.0)] if b else [R(" ",size=6)]],align=PP_ALIGN.CENTER)
    if i<len(chain)-1:
        text(s,x+bw+0.01,6.18,0.34,0.5,[[R("←",size=20,color=RED,bold=True,font=LAT_FONT)]],align=PP_ALIGN.CENTER)

# ============================================================ 8 BLACK BOX
s=slide(); bg(s); footer(s,8)
kicker(s, 10.6, 0.7, "الصندوق الأسود · 10 دقائق", w=2.5, color=AMBER)
heading(s, 0.6, 1.35, 12.1, [("أسئلة الفورم — ",WHITE),("إحنا بنسمع فعلًا",RED)], size=32)
bb=[
 ("🕵  من غير أسماء","نقرا أكتر 3–4 أسئلة تكرارًا ونجاوب عليها مباشرة — الخصوصية محفوظة، والجرأة مكافأة."),
 ("⚡  رد سريع ومباشر","من غير لف ولا دفاعية. سؤال واضح يستحق إجابة واضحة."),
 ("⏰  ركّز على سؤال واحد","أي سؤال عن الخوف من إدارة الوقت مع الدراسة ياخد وقت أكبر — لأنه بيفتح الفقرة الجاية مباشرة."),
]
x=0.6
for (t,d) in bb:
    card(s,x,2.6,3.94,2.3,accent=True)
    text(s,x+0.28,2.9,3.4,0.7,[[R(t,size=16,color=WHITE,bold=True,line=1.2)]])
    text(s,x+0.28,3.75,3.45,1.0,[[R(d,size=13.5,color=GREY,line=1.4)]])
    x+=4.09
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,1.6,5.4,10.1,1.1,fill=PANEL,line=LINE,radius=0.1)
text(s,1.95,5.62,9.4,0.7,[[R("الرسالة الضمنية: «مش بس بنتكلم فيكم… إحنا بنقرا اللي بتكتبوه، وأسئلتكم بتشكّل الميتنج.»",size=15,color=WHITE,bold=True,line=1.4)]],align=PP_ALIGN.CENTER)

# ============================================================ 9 MINDSET
s=slide(); bg(s); footer(s,9)
kicker(s, 8.0, 0.55, "الفقرة الجوهرية · 30 دقيقة — 3 رسائل بعمق، مش عشرة سطحيين", w=5.2, color=GREEN)
heading(s, 0.6, 1.1, 12.1, [("أولًا: الـ ",WHITE),("Mindset",GREEN),("  · 8 دقائق",WHITE)], size=30)
card(s,0.6,2.1,5.9,3.9, fill=RGBColor(0x0E,0x1C,0x16), line=RGBColor(0x14,0x4A,0x36))
text(s,0.9,2.4,5.3,0.6,[[R("🧊  الفكرة المحورية: افصل المشاعر عن الالتزام",size=16,color=GREEN,bold=True,line=1.2)]])
text(s,0.9,3.25,5.35,1.6,[[R("«ضغطك، مزاجك، ظروفك — دي حاجات إنسانية وطبيعية ومحترمة. لكن التزامك وكلمتك قدام التيم… حاجة تانية خالص.»",size=15,color=WHITE,bold=True,line=1.55)]])
text(s,0.9,5.05,5.35,0.9,[[R("المحترف الحقيقي مش اللي معندوش ضغط — هو اللي ",size=14,color=GREY),R("بيسلّم حتى وهو تحت ضغط",size=14,color=WHITE,bold=True),R(".",size=14,color=GREY,line=1.4)]])
card(s,6.85,2.1,5.9,3.9,accent=True)
text(s,7.15,2.4,5.3,0.6,[[R("🥊  الفكرة التانية: الفعل بدل الكلام",size=16,color=RED,bold=True)]])
text(s,7.15,3.25,5.35,1.0,[[R("«متقولش هتعمل إيه… اعمله، وخلّي النتيجة تتكلم.»",size=16,color=WHITE,bold=True,line=1.5)]])
text(s,7.15,4.5,5.35,1.2,[[R("الكلام الكتير عن النية بيستهلك ",size=14,color=GREY),R("نفس الطاقة",size=14,color=WHITE,bold=True),R(" اللي المفروض تتصرف في التنفيذ.",size=14,color=GREY,line=1.45)]])

# ============================================================ 10 PARETO
s=slide(); bg(s); footer(s,10)
kicker(s, 10.5, 0.55, "تانيًا: مهارة الوقت · 10 دقائق", w=2.6)
text(s,0.6,1.35,4.2,1.6,[[R("80/20",size=96,color=RED,bold=True,font=LAT_FONT,line=0.95)]],align=PP_ALIGN.CENTER)
heading(s, 4.6, 1.5, 8.1, [("قاعدة ",WHITE),("پاريتو",RED),(" — ركّز في الـ 20% اللي بتفرق",WHITE)], size=27)
text(s,4.6,2.75,8.1,1.0,[[R("80% من نتيجتك بتيجي من ",size=16,color=GREY),R("20% بس",size=16,color=RED,bold=True),
  R(" من مجهودك. المشكلة إن أغلبنا بيوزّع وقته بالتساوي على كل حاجة، فيضيع وقت في تفاصيل مالهاش تأثير.",size=16,color=GREY,line=1.5)]])
par=[
 ("🎯  قبل أي تاسك","اسأل: «إيه الجزء اللي لو خلصته… هيخلي 80% من التاسك خلص؟» وابدأ بيه."),
 ("📚  في المذاكرة","20% من المنهج (المفاهيم الأساسية) بتغطّي 80% من فهمك للباقي — ركّز عليها الأول."),
 ("🌅  أول ساعة في يومك","أعلى طاقة ذهنية تروح للمهمة الأهم — مش للإيميلات والتنظيف السطحي."),
]
x=0.6
for (t,d) in par:
    card(s,x,4.25,3.94,2.2,accent=True)
    text(s,x+0.28,4.5,3.4,0.6,[[R(t,size=15.5,color=WHITE,bold=True)]])
    text(s,x+0.28,5.2,3.45,1.1,[[R(d,size=13.5,color=GREY,line=1.45)]])
    x+=4.09

# ============================================================ 11 12 WEEK YEAR
s=slide(); bg(s); footer(s,11)
kicker(s, 9.7, 0.55, "تطبيق إيقاع التيم — من المستند المكتوب", w=3.4, color=AMBER)
heading(s, 0.6, 1.1, 12.1, [("الـ ",WHITE),("12 Week Year",RED),(" — السنة مبقتش 12 شهر",WHITE)], size=30)
card(s,0.6,2.2,5.9,2.3,fill=PANEL,line=LINE)
text(s,0.9,2.45,5.3,0.5,[[R("🐢  ليه الأهداف السنوية بتفشل؟",size=16,color=WHITE,bold=True)]])
text(s,0.9,3.15,5.35,1.2,[[R("التخطيط السنوي بيخلي الدماغ يأجّل لآخر 3 شهور — الإحساس بالإلحاح بيموت، وكل حاجة «بكرة».",size=14.5,color=GREY,line=1.5)]])
card(s,6.85,2.2,5.9,2.3,accent=True)
text(s,7.15,2.45,5.3,0.5,[[R("🔥  الحل: كل 12 أسبوع = سنة كاملة",size=16,color=RED,bold=True,line=1.2)]])
text(s,7.15,3.25,5.35,1.2,[[R("اتعامل مع كل 12 أسبوع كأنها سنة مستقلة، بهدف واضح وخطة أسبوعية — الإلحاح بيفضل شغّال طول الوقت.",size=14.5,color=GREY,line=1.5)]])
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,0.6,4.85,12.13,1.5,fill=PANEL2,line=LINE_R,line_w=1.25,radius=0.07)
text(s,0.95,5.15,11.5,1.1,[[R("🎯 تطبيق Phantoms: ",size=15.5,color=RED,bold=True),
  R("كل ربع سنة (≈ ترم دراسي) ليه هدف تقني واحد واضح ومحدد وقابل للقياس — مثلاً «مشروع كامل جاهز» أو «مهارة معينة متقَنة 100%» — وكل أسبوع فيه تسليم صغير بيقرّب من الهدف.",size=15.5,color=WHITE,line=1.55)]])

# ============================================================ 12 COMMUNICATION
s=slide(); bg(s); footer(s,12)
kicker(s, 10.35, 0.55, "تالتًا: مهارة التواصل · 12 دقيقة", w=2.75, color=AMBER)
heading(s, 0.6, 1.05, 12.1, [("الحضور الواثق — ",WHITE),("متتكلمش زي الموظف",RED)], size=29)
card(s,0.6,1.95,5.9,2.85, fill=RGBColor(0x24,0x10,0x10), line=RGBColor(0x6E,0x14,0x10))
text(s,0.9,2.2,5.3,0.5,[[R("🙊  الموظف — اللي بيتنسى",size=16,color=RED,bold=True)]])
bullets(s,0.9,2.85,5.35,1.8,[
 [R("«أنا عملت اللي طلبتوه.»",size=14,color=WHITE,bold=True)],
 [R("بيشتغل بصمت، يوافق أول واحد، بيتكلم آخره من غير رأي، ويعتذر قبل ما يقول رأيه.",size=14,color=GREY)],
],gap=8)
card(s,6.85,1.95,5.9,2.85, fill=RGBColor(0x0E,0x1C,0x16), line=RGBColor(0x14,0x4A,0x36))
text(s,7.15,2.2,5.3,0.5,[[R("🦅  صاحب المسؤولية — اللي بيتذكر",size=16,color=GREEN,bold=True)]])
bullets(s,7.15,2.85,5.35,1.8,[
 [R("«عملت X، ولاحظت Y، فعملت Z كمان.»",size=14,color=WHITE,bold=True)],
 [R("عنده موقف واضح: «أنا شايف الحل التاني أفضل، وده السبب» — بهدوء وثبات، من غير تلعثم.",size=14,color=GREY)],
],gap=8,marker_color=GREEN)
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,0.6,5.05,5.9,1.55,fill=PANEL2,line=LINE_R,line_w=1,radius=0.08)
text(s,0.9,5.28,5.35,1.2,[[R("💪 الثقة مش عناد ولا رفع صوت — إنك تقول رأيك بوضوح وهدوء، من غير ما تعتذر عن رأيك قبل ما تقوله. الشغل الكويس ",size=13.5,color=GREY,line=1.4),R("مش",size=13.5,color=RED,bold=True),R(" بيتكلم عن نفسه لوحده.",size=13.5,color=GREY)]])
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,6.85,5.05,5.9,1.55,fill=PANEL,line=LINE,line_w=1,radius=0.08)
text(s,7.15,5.28,5.35,1.2,[[R("🧪 تمرين يتنفّذ من بكرة: ",size=13.5,color=AMBER,bold=True),R("قبل أي اجتماع حضّر ",size=13.5,color=GREY),R("جملة واحدة واضحة",size=13.5,color=WHITE,bold=True),R(" تعبّر بيها عن رأيك. متدخلش فاضي تستني — تيجي مجهّز بحلول مش أسئلة، وبأرقام مش انطباعات.",size=13.5,color=GREY,line=1.4)]])

# ============================================================ 13 BOOKS
s=slide(); bg(s); footer(s,13)
kicker(s, 9.9, 0.55, "30 ثانية لكل كتاب — التفاصيل في المستند المكتوب", w=3.2)
heading(s, 0.6, 1.05, 12.1, [("مكتبة ",WHITE),("Phantoms",RED),(" — العادات والعقل والتركيز",WHITE)], size=29)
books=[
 ("atomic-habits.jpg","Atomic Habits","James Clear","بناء عادة التطوّر بزيادات صغيرة يومية من غير ضغط — 1% كل يوم."),
 ("deep-work.jpg","Deep Work","Cal Newport","ساعتين تركيز عميق بتنتجوا أكتر من أسبوع كامل من التشتّت."),
 ("12-week-year.png","The 12 Week Year","B. Moran & M. Lennington","التنفيذ في دورات قصيرة عالية الكثافة — بدل التأجيل السنوي."),
 ("daily-stoic.jpg","The Daily Stoic","Ryan Holiday","ثبات انفعالي تحت الضغط — حكمة يومية للسنة كاملة."),
]
x=0.6
for (fn,t,au,d) in books:
    card(s,x,1.95,2.94,4.5,fill=PANEL,line=LINE)
    add_pic_crop(s, os.path.join(A,"books",fn), x+0.22,2.18,2.5,2.55, line=LINE_R)
    text(s,x+0.22,4.95,2.5,0.7,[[R(t,size=15,color=WHITE,bold=True,font=LAT_FONT,line=1.1)]])
    text(s,x+0.22,5.62,2.5,0.3,[[R(au,size=11,color=RED,bold=True,font=LAT_FONT)]])
    text(s,x+0.22,5.95,2.55,0.85,[[R(d,size=11.5,color=GREY,line=1.3)]])
    x+=3.06

# ============================================================ 14 DR BRAVO 1
s=slide(); bg(s); footer(s,14)
s.shapes.add_picture(os.path.join(A,"bravo-cutout.png"), Inches(0.3), Inches(2.2), height=Inches(4.7))
kicker(s, 8.9, 0.6, "📌 تحليل كامل — ليه بنتعلم منه بالظبط؟", w=4.2, color=AMBER)
heading(s, 4.4, 1.2, 8.3, [("د. محمود برافو  ",WHITE),("Dr. Mahmoud Bravo",RED)], size=28)
text(s,4.4,2.15,8.3,1.0,[[R("مدرب شركات دولي وكوتش أداء، متحدث رئيسي (Keynote)، وخبير تسويق وتواصل في الرعاية الصحية — والأهم: بنى مهارات العرض والحضور دي من مشوار تنفيذي حقيقي، مش من كتب.",size=13.5,color=GREY,line=1.45)]])
stats=[("1500+","مشروع عبر MENA"),("3000+","متخصص درّبهم"),("50+","حدث كمتحدث رئيسي"),("30+","جائزة عالمية وإقليمية")]
x=4.4
for (v,l) in stats:
    card(s,x,3.25,1.98,1.0,fill=PANEL,line=LINE)
    text(s,x,3.36,1.98,0.5,[[R(v,size=22,color=RED,bold=True,font=LAT_FONT)]],align=PP_ALIGN.CENTER)
    text(s,x,3.86,1.98,0.35,[[R(l,size=10.5,color=GREY)]],align=PP_ALIGN.CENTER)
    x+=2.1
tl=[("2010","صيدلة عين شمس + دبلوم تسويق AUC → Novartis"),
    ("2013–16","Brand Manager: نمو 35% فوق التارجت + MBA جنيف"),
    ("2017–19","Abbott Nutrition ثم Medtronic — تسويق إقليمي في 9 دول"),
    ("2021–24","MD لـ Healthcare COMs · DBA · رئيس تميز الأعمال بـ Utopia")]
x=4.4
for (yr,d) in tl:
    card(s,x,4.5,1.98,1.55,fill=PANEL2,line=LINE_R)
    text(s,x+0.12,4.66,1.74,0.4,[[R(yr,size=13,color=RED,bold=True,font=LAT_FONT)]])
    text(s,x+0.12,5.1,1.78,0.9,[[R(d,size=10.5,color=GREY,line=1.25)]])
    x+=2.1
text(s,4.4,6.3,8.3,0.5,[[R("المصادر: mahmoudbravo.com · صفحة About · @drmahmoudbravo · TEDxEJUST",size=10.5,color=DIM,font=LAT_FONT)]])

# ============================================================ 15 DR BRAVO 2
s=slide(); bg(s); footer(s,15)
kicker(s, 9.3, 0.55, "تحليل الأسلوب — إيه اللي بيخليه مؤثر؟", w=3.8, color=AMBER)
heading(s, 0.6, 1.05, 12.1, [("5 دروس من برافو… ",WHITE),("بنطبّقها في ميتنج النهاردة",RED)], size=27)
lessons=[
 ("1 · المزج بين الاستراتيجي والإنساني","موقعه بينصّ عليها: عمق استراتيجي بتواضع وتعاطف. الدرس: الحزم في الرسالة، الدفء في التوصيل — ده توازن «البرمشن» و«التحوّل لصحاب» في خطتنا."),
 ("2 · المصداقية بالأرقام مش بالكلام","1500+ مشروع، 3000+ متدرب، 9 دول. الدرس: في أول جملتين للتعريف — اضرب بإنجاز مختصر بأرقام واقفل بسرعة. الأرقام بتتكلم."),
 ("3 · الصوت أول سلاح","محتواه كله عن الصوت: «Don't Underestimate Your Voice» و«Find Your Authority». الدرس: أول 30 ثانية افتتح بصوت ثابت ومرتاح — الانطباع الأول سمعي."),
 ("4 · التأثير الهادئ بدل البيع العنيف","أشهر حواراته: «كيف تبيع دون أن تبيع — فن التأثير الهادئ». الدرس: القيادة بالإقناع الهادي أقوى من التهديد — الـ Strikes بهدوء وحزم."),
]
x0,y0=0.6,1.9; cw,ch=5.98,1.75
for i,(t,d) in enumerate(lessons):
    r,c=divmod(i,2)
    x=x0+c*(cw+0.17); y=y0+r*(ch+0.18)
    card(s,x,y,cw,ch,accent=True)
    text(s,x+0.3,y+0.2,cw-0.55,0.55,[[R(t,size=14.5,color=WHITE,bold=True,line=1.15)]])
    text(s,x+0.3,y+0.82,cw-0.6,0.85,[[R(d,size=12.5,color=GREY,line=1.3)]])
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,0.6,5.75,12.13,1.05,fill=PANEL2,line=LINE_R,line_w=1.25,radius=0.08)
text(s,0.95,5.95,11.5,0.8,[[R("👂 5 · الاستماع النشط: ",size=14,color=RED,bold=True),
  R("حلقاته عن Active Listening مع المشاعر — ده اللي بيخلي «الصندوق الأسود» مش فقرة شكلية. وقانونه: «الخبرة التقنية بتفتحلك الباب… لكن المهارات الناعمة هي اللي بتحدّد توصل لفين.»",size=13.5,color=WHITE,line=1.4)]])

# ============================================================ 16 POST DOC
s=slide(); bg(s); footer(s,16)
kicker(s, 9.0, 0.55, "المستند المكتوب — يترسل بعد الميتنج بساعات", w=4.1, color=AMBER)
heading(s, 0.6, 1.05, 12.1, [("اللي هنقوله ",WHITE),("بعمق في الجروب",RED),("… بدل ما يضيع في الزحمة",WHITE)], size=27)
docs=[
 ("📅  الـ 12 Week Year خطوة بخطوة","الفكرة كاملة + التطبيق على التيم: هدف ربعي محدد، خطة أسبوعية، تسليم أسبوعي صغير."),
 ("🧭  إطار اختيار فكرة المشروع","Product-Market Fit مبسّط: ابدأ من مشكلة حقيقية → مين المستخدم؟ → ابنِ أبسط نسخة (MVP) وجرّبها بسرعة."),
 ("📚  ملخصات أطول للكتب الأربعة","Atomic Habits · Deep Work · 12 Week Year · The Daily Stoic — نقاط قابلة للتطبيق مش مراجعات."),
 ("🗺️  الـ Roadmap الكامل بالتواريخ","خطة الترم الأول لكل لجنة بتواريخها — مرجع دائم يرجعوله في أي وقت."),
]
x0,y0=0.6,1.95
for i,(t,d) in enumerate(docs):
    r,c=divmod(i,2)
    x=x0+c*(5.98+0.17); y=y0+r*(1.7+0.2)
    card(s,x,y,5.98,1.7,fill=PANEL,line=LINE)
    text(s,x+0.3,y+0.25,5.4,0.5,[[R(t,size=15.5,color=WHITE,bold=True)]])
    text(s,x+0.3,y+0.85,5.45,0.8,[[R(d,size=13,color=GREY,line=1.4)]])
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,0.6,5.95,12.13,0.95,fill=PANEL2,line=LINE_R,line_w=1,radius=0.1)
text(s,0.95,6.15,11.5,0.6,[[R("⏱ لو الوقت ضاق: ",size=13.5,color=AMBER,bold=True),
  R("اختصر عروض اللجان لـ 10 دقائق بدل 13 — ده يوفّر 9 دقايق تروح لفقرة الـ Mindset. المكتوب يفضل هو شبكة الأمان.",size=13.5,color=WHITE,line=1.35)]])

# ============================================================ 17 CLOSING
s=slide(); bg(s); footer(s,17)
kicker(s, 11.0, 0.7, "الختام · 5 دقائق", w=2.1, color=AMBER)
text(s,0.6,1.45,12.1,0.6,[[R("الخلاصة في جملة واحدة",size=30,color=WHITE,bold=True)]],align=PP_ALIGN.CENTER)
text(s,1.6,2.35,10.1,1.5,[
  [R("«النهاردة مش بس سمعتوا خطة…",size=26,color=WHITE,bold=True,after=6,line=1.35)],
  [R("انتوا بقيتوا جزء من ",size=26,color=WHITE,bold=True),R("كيان بيصنع فرق",size=26,color=RED,bold=True),R(".»",size=26,color=WHITE,bold=True,line=1.35)],
],align=PP_ALIGN.CENTER)
cl=[("🙏  قدّر الحضور","كل واحد يحس إنه اتقدّر لالتزامه — مش إنه اتهدّد بس."),
    ("🔁  كرّر الجملة التأسيسية","الشعار اللي افتتحت بيه يتقفل بيه — بيثبت في الذاكرة."),
    ("🔻  انزل من المنصة","ابتسامة، تشمير كم، وجلوس وسط التيم — الإشارة الجسدية أقوى من كلام.")]
x=0.6
for (t,d) in cl:
    card(s,x,4.4,3.94,2.0,accent=True)
    text(s,x+0.28,4.65,3.4,0.6,[[R(t,size=15.5,color=WHITE,bold=True)]])
    text(s,x+0.28,5.35,3.45,1.0,[[R(d,size=13.5,color=GREY,line=1.45)]])
    x+=4.09

# ============================================================ 18 FINALE
s=slide(); bg(s)
logo(s, 5.35, 0.75, 2.6)
shape(s, MSO_SHAPE.RECTANGLE, 3.4, 3.55, 6.53, 0.025, fill=RED)
text(s,1.4,3.75,10.5,1.7,[
  [R("«لحد هنا… كلام الشغل والرسميات خلص.",size=26,color=WHITE,bold=True,after=8,line=1.4)],
  [R("من اللحظة دي إحنا ",size=26,color=WHITE,bold=True),R("إخوات وصحاب",size=26,color=RED,bold=True),R(".»",size=26,color=WHITE,bold=True,line=1.4)],
],align=PP_ALIGN.CENTER)
text(s,2.4,5.55,8.5,0.9,[
  [R("يلا نفصّل… ونشوف مين هيسد ومين هيخسر في اللعب. 🎮",size=16,color=GREY,bold=True,after=6,line=1.4)],
  [R("اللعب يبدأ فورًا — بلا فاصل زمني. الانتقال المفاجئ نفسه هو اللي بيكسر التوتر.",size=13.5,color=AMBER,bold=True,line=1.35)],
],align=PP_ALIGN.CENTER)
shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,4.5,6.7,4.33,0.55,fill=PANEL2,line=LINE_R,line_w=1,radius=0.5)
text(s,4.5,6.78,4.33,0.4,[[R("👻  Phantoms — السنة دي بتتفرق",size=14,color=WHITE,bold=True)]],align=PP_ALIGN.CENTER)

out=os.path.join(HERE,"Phantoms-Meeting-01.pptx")
prs.save(out)
print("SAVED", out, os.path.getsize(out), "bytes ;", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
